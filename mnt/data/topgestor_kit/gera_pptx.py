# -*- coding: utf-8 -*-
"""
Gera um PowerPoint com:
- Slide de título
- Tabela de performance (a partir de outputs/tabela_performance.xlsx)
- Slides com os gráficos PNG gerados em outputs/
Requer: python-pptx, pandas
"""
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

BASE = Path(__file__).resolve().parent
OUT = BASE / "outputs"

def add_title_slide(prs, titulo, subtitulo):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = titulo
    slide.placeholders[1].text = subtitulo
    return slide

def add_table_slide(prs, title, df: pd.DataFrame):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    rows, cols = df.shape[0] + 1, df.shape[1]
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # header
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = str(col)

    # body
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            val = df.iat[i, j]
            if isinstance(val, float):
                # Try to detect % columns by header
                hdr = df.columns[j].lower()
                if "vol" in hdr or "cdi" in hdr or "fundo" in hdr or "cdi" == hdr:
                    table.cell(i+1, j).text = f"{val:.2%}"
                else:
                    table.cell(i+1, j).text = f"{val:.4f}" if abs(val) < 1 else f"{val:,.2f}"
            else:
                table.cell(i+1, j).text = str(val)

    # basic style: align center header
    for j in range(cols):
        cell = table.cell(0, j)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.bold = True
                p.runs[0].font.size = Pt(12)

    return slide

def add_image_slide(prs, title, image_path: Path):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # center image
    pic = slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.5), height=Inches(5.0))
    return slide

def main():
    prs = Presentation()
    add_title_slide(prs, "Carta Mensal – Engenheiras da Bolsa", "Gráficos e Tabelas")

    # Tabela
    xlsx = OUT / "tabela_performance.xlsx"
    if xlsx.exists():
        df = pd.read_excel(xlsx, sheet_name="Performance")
        add_table_slide(prs, "Tabela de Performance", df)
    else:
        print("Aviso: arquivo outputs/tabela_performance.xlsx não encontrado. Pulei a tabela.")

    # Gráficos: adiciona cada PNG em um slide
    for png in sorted(OUT.glob("*.png")):
        title = png.stem.replace("_", " ").title()
        add_image_slide(prs, title, png)

    out_pptx = OUT / "carta_slides.pptx"
    prs.save(out_pptx)
    print("Gerado:", out_pptx)

if __name__ == "__main__":
    main()
